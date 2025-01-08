import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
import pptx
import pandas as pd
import os
import time
import torch
from transformers import pipeline, AutoTokenizer, AutoModelForSeq2SeqLM
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np

st.set_page_config(
    page_title="SDLC Automate APP",
    page_icon="images/favicon.png"
)

# Hide Streamlit branding, menu, and footer
hide_streamlit_style = """
    <style>
    #MainMenu {visibility: hidden;} /* Hides the "Manage app" menu */
    footer {visibility: hidden;} /* Hides the Streamlit footer */
    header {visibility: hidden;} /* Hides the header */
    ._link_gzau3_10 {display: none;} /* Hides "Hosted with Streamlit" */
    .stDeployButton {display: none !important;} /* Hides the deploy button */
    </style>
"""
st.markdown(hide_streamlit_style, unsafe_allow_html=True)

# Load Hugging Face model for text generation
@st.cache_resource
def load_model():
    model_name = "facebook/bart-large-cnn"
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
    return tokenizer, model

tokenizer, model = load_model()

# Load Sentence Transformer for embedding-based retrieval
@st.cache_resource
def load_embedding_model():
    return SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

embedding_model = load_embedding_model()

# Streamlit sidebar setup
with st.sidebar:
    st.title("Your BRD Documents")
    uploaded_file = st.file_uploader("Upload a file to generate user stories", type=["pdf", "docx", "txt", "xlsx", "pptx"])

# Function to extract text from various file types
@st.cache_resource
def extract_text_from_file(file):
    text = ""
    file_ext = os.path.splitext(file.name)[1].lower()

    if file_ext == ".pdf":
        pdf_reader = PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()

    elif file_ext == ".docx":
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif file_ext == ".txt":
        text = file.read().decode("utf-8")

    elif file_ext in [".xlsx", ".xls"]:
        df = pd.read_excel(file)
        text = df.to_string()

    elif file_ext in [".pptx", ".ppt"]:
        ppt = pptx.Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    return text

# Process uploaded file and extract text
@st.cache_resource
def process_uploaded_file(uploaded_file):
    return extract_text_from_file(uploaded_file) if uploaded_file else ""

# Function to create FAISS vector store from extracted text
@st.cache_resource
def create_vector_store(text):
    text_chunks = text.split("\n")
    embeddings = embedding_model.encode(text_chunks)
    
    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(np.array(embeddings).astype("float32"))

    return text_chunks, index

# Function to search in vector store
def retrieve_similar_text(query, text_chunks, index, k=3):
    query_embedding = embedding_model.encode([query])
    distances, indices = index.search(np.array(query_embedding).astype("float32"), k)
    return [text_chunks[i] for i in indices[0]]

# Function to generate text using Hugging Face model
def generate_text(prompt):
    input_ids = tokenizer(prompt, return_tensors="pt", truncation=True, padding="max_length", max_length=512).input_ids
    output = model.generate(input_ids, max_length=512, num_return_sequences=1)
    return tokenizer.decode(output[0], skip_special_tokens=True)

# Streamlit app setup
st.header("BRD to User Story, Test Case, Cucumber Script, and Selenium Script")

# Set up tabs for different functionalities
tab1, tab2, tab3, tab4 = st.tabs(["User Story Generation", "User Story to Test Case", "Test Case to Cucumber Script", "Test Case to Selenium Script"])

# User Story Generation Tab
with tab1:
    start_time = time.time()
    if uploaded_file:
        text = process_uploaded_file(uploaded_file)
        if text:
            text_chunks, index = create_vector_store(text)
            prompt_message = (
                "Think of yourself as a senior business analyst. Your responsibility is to read the Business Requirement Document "
                "and write the User Stories according to that BRD. Think step-by-step and write all possible user stories."
            )
            matches = retrieve_similar_text(prompt_message, text_chunks, index, k=3)
            response = generate_text(" ".join(matches))
            st.write(response)
            st.write(f"Processing time: {time.time() - start_time:.2f} seconds")
    else:
        st.write("Please upload a BRD document.")

# User Story to Test Case Tab
with tab2:
    st.subheader("Convert User Story to Test Case")
    user_story_text = st.text_area("Enter the user story:")

    if st.button("Generate Test Cases"):
        if user_story_text:
            prompt = "Generate test cases for the following user story:\n" + user_story_text
            response = generate_text(prompt)
            st.write(response)
        else:
            st.write("Please enter a user story.")

# Test Case to Cucumber Script Tab
with tab3:
    st.subheader("Convert Test Case to Cucumber Script")
    test_case_text = st.text_area("Enter the test case:")

    if st.button("Generate Cucumber Script"):
        if test_case_text:
            prompt = "Convert this test case into a Cucumber script using Gherkin syntax:\n" + test_case_text
            response = generate_text(prompt)
            st.write(response)
        else:
            st.write("Please enter a test case.")

# Test Case to Selenium Script Tab
with tab4:
    st.subheader("Convert Test Case to Selenium Script")
    selenium_test_case_text = st.text_area("Enter the test case:")

    if st.button("Generate Selenium Script"):
        if selenium_test_case_text:
            prompt = "Convert this test case into a Selenium WebDriver script using Python:\n" + selenium_test_case_text
            response = generate_text(prompt)
            st.write(response)
        else:
            st.write("Please enter a test case.")
